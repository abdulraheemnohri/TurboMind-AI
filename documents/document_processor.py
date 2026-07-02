#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
TurboMind AI - Document Processor
=================================
Processes various document formats (PDF, DOCX, TXT, etc.)
"""

import os
import re
from pathlib import Path
from typing import Dict, Any, List, Optional, Union
from dataclasses import dataclass, field
from enum import Enum
import json


class DocumentType(Enum):
    """Supported document types"""
    PDF = "pdf"
    DOCX = "docx"
    TXT = "txt"
    MARKDOWN = "md"
    HTML = "html"
    CSV = "csv"
    EXCEL = "xlsx"
    POWERPOINT = "pptx"
    UNKNOWN = "unknown"


@dataclass
class DocumentMetadata:
    """Metadata about a document"""
    file_path: str
    file_name: str
    file_type: DocumentType
    file_size: int  # bytes
    page_count: int = 0
    word_count: int = 0
    char_count: int = 0
    created_at: Optional[float] = None
    modified_at: Optional[float] = None
    title: str = ""
    author: str = ""
    language: str = "en"


@dataclass
class DocumentChunk:
    """A chunk of document content"""
    text: str
    chunk_id: int
    start_page: int = 0
    end_page: int = 0
    metadata: Dict[str, Any] = field(default_factory=dict)


@dataclass
class ProcessedDocument:
    """A fully processed document"""
    metadata: DocumentMetadata
    content: str
    chunks: List[DocumentChunk] = field(default_factory=list)
    sections: List[Dict[str, Any]] = field(default_factory=list)
    embeddings: Optional[List[List[float]]] = None


class DocumentProcessor:
    """
    Processes various document formats and extracts text content.
    Handles PDF, DOCX, TXT, Markdown, HTML, CSV, Excel, and PowerPoint files.
    """
    
    # Chunk configuration
    DEFAULT_CHUNK_SIZE = 1000  # tokens
    DEFAULT_CHUNK_OVERLAP = 200  # tokens
    
    def __init__(self, chunk_size: int = DEFAULT_CHUNK_SIZE, chunk_overlap: int = DEFAULT_CHUNK_OVERLAP):
        """
        Initialize the document processor.
        
        Args:
            chunk_size: Size of each chunk in tokens
            chunk_overlap: Overlap between chunks in tokens
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self._temp_files = []  # Track temp files for cleanup
        
        print(f"📄 Document Processor initialized (chunk: {chunk_size}, overlap: {chunk_overlap})")
    
    def detect_document_type(self, file_path: str) -> DocumentType:
        """Detect document type from file extension"""
        ext = Path(file_path).suffix.lower().replace('.', '')
        
        type_map = {
            'pdf': DocumentType.PDF,
            'docx': DocumentType.DOCX,
            'txt': DocumentType.TXT,
            'md': DocumentType.MARKDOWN,
            'markdown': DocumentType.MARKDOWN,
            'html': DocumentType.HTML,
            'htm': DocumentType.HTML,
            'csv': DocumentType.CSV,
            'xlsx': DocumentType.EXCEL,
            'xls': DocumentType.EXCEL,
            'pptx': DocumentType.POWERPOINT,
            'ppt': DocumentType.POWERPOINT,
        }
        
        return type_map.get(ext, DocumentType.UNKNOWN)
    
    def get_metadata(self, file_path: str) -> DocumentMetadata:
        """Get metadata from a document file"""
        file_path_obj = Path(file_path)
        file_type = self.detect_document_type(file_path)
        
        stat = file_path_obj.stat()
        
        metadata = DocumentMetadata(
            file_path=str(file_path_obj),
            file_name=file_path_obj.name,
            file_type=file_type,
            file_size=stat.st_size,
            created_at=stat.st_ctime,
            modified_at=stat.st_mtime
        )
        
        # Try to extract additional metadata based on type
        if file_type == DocumentType.PDF:
            metadata = self._get_pdf_metadata(file_path_obj, metadata)
        elif file_type == DocumentType.DOCX:
            metadata = self._get_docx_metadata(file_path_obj, metadata)
        
        return metadata
    
    def _get_pdf_metadata(self, file_path: Path, metadata: DocumentMetadata) -> DocumentMetadata:
        """Extract metadata from PDF file"""
        try:
            # Try to use PyPDF2
            try:
                from PyPDF2 import PdfReader
                with open(file_path, 'rb') as f:
                    reader = PdfReader(f)
                    metadata.title = reader.metadata.get('/Title', '') or metadata.title
                    metadata.author = reader.metadata.get('/Author', '') or metadata.author
                    metadata.page_count = len(reader.pages)
            except ImportError:
                pass
        except Exception:
            pass
        
        return metadata
    
    def _get_docx_metadata(self, file_path: Path, metadata: DocumentMetadata) -> DocumentMetadata:
        """Extract metadata from DOCX file"""
        try:
            from docx import Document
            doc = Document(file_path)
            metadata.page_count = len(doc.paragraphs)  # Approximate
        except ImportError:
            pass
        except Exception:
            pass
        
        return metadata
    
    def process(self, file_path: str, **kwargs) -> Optional[ProcessedDocument]:
        """
        Process a document file and extract content.
        
        Args:
            file_path: Path to the document file
            **kwargs: Additional processing options
            
        Returns:
            ProcessedDocument with extracted content and metadata
        """
        file_path_obj = Path(file_path)
        
        if not file_path_obj.exists():
            print(f"❌ File not found: {file_path}")
            return None
        
        # Get metadata
        metadata = self.get_metadata(file_path)
        
        # Extract content based on type
        content = ""
        
        try:
            if metadata.file_type == DocumentType.PDF:
                content = self._extract_pdf(file_path_obj)
            elif metadata.file_type == DocumentType.DOCX:
                content = self._extract_docx(file_path_obj)
            elif metadata.file_type == DocumentType.TXT:
                content = self._extract_txt(file_path_obj)
            elif metadata.file_type == DocumentType.MARKDOWN:
                content = self._extract_markdown(file_path_obj)
            elif metadata.file_type == DocumentType.HTML:
                content = self._extract_html(file_path_obj)
            elif metadata.file_type == DocumentType.CSV:
                content = self._extract_csv(file_path_obj)
            elif metadata.file_type == DocumentType.EXCEL:
                content = self._extract_excel(file_path_obj)
            elif metadata.file_type == DocumentType.POWERPOINT:
                content = self._extract_powerpoint(file_path_obj)
            else:
                print(f"⚠️  Unsupported file type: {metadata.file_type}")
                return None
        except Exception as e:
            print(f"❌ Error processing {file_path}: {e}")
            return None
        
        # Clean and normalize content
        content = self._clean_content(content)
        
        # Update word and character counts
        metadata.word_count = len(content.split())
        metadata.char_count = len(content)
        
        # Create chunks
        chunks = self._create_chunks(content, metadata.file_type)
        
        # Create processed document
        processed = ProcessedDocument(
            metadata=metadata,
            content=content,
            chunks=chunks
        )
        
        print(f"✅ Processed: {file_path_obj.name} ({len(content)} chars, {len(chunks)} chunks)")
        return processed
    
    def _extract_pdf(self, file_path: Path) -> str:
        """Extract text from PDF file"""
        content = ""
        
        # Try PyPDF2
        try:
            from PyPDF2 import PdfReader
            with open(file_path, 'rb') as f:
                reader = PdfReader(f)
                for page in reader.pages:
                    content += page.extract_text() + "\n"
        except ImportError:
            pass
        
        # Try pdfplumber (better for complex PDFs)
        if not content.strip():
            try:
                import pdfplumber
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        content += page.extract_text() + "\n"
            except ImportError:
                pass
        
        return content
    
    def _extract_docx(self, file_path: Path) -> str:
        """Extract text from DOCX file"""
        content = ""
        
        try:
            from docx import Document
            doc = Document(file_path)
            for para in doc.paragraphs:
                content += para.text + "\n"
        except ImportError:
            pass
        
        return content
    
    def _extract_txt(self, file_path: Path) -> str:
        """Extract text from TXT file"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    
    def _extract_markdown(self, file_path: Path) -> str:
        """Extract text from Markdown file"""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    
    def _extract_html(self, file_path: Path) -> str:
        """Extract text from HTML file"""
        try:
            from bs4 import BeautifulSoup
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
                return soup.get_text()
        except ImportError:
            # Fallback: just read as text
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
    
    def _extract_csv(self, file_path: Path) -> str:
        """Extract text from CSV file"""
        import csv
        content = ""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            reader = csv.reader(f)
            for row in reader:
                content += ", ".join(row) + "\n"
        return content
    
    def _extract_excel(self, file_path: Path) -> str:
        """Extract text from Excel file"""
        content = ""
        
        try:
            import pandas as pd
            df = pd.read_excel(file_path)
            content = df.to_string()
        except ImportError:
            pass
        
        return content
    
    def _extract_powerpoint(self, file_path: Path) -> str:
        """Extract text from PowerPoint file"""
        content = ""
        
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
        except ImportError:
            pass
        
        return content
    
    def _clean_content(self, content: str) -> str:
        """Clean and normalize extracted content"""
        # Remove excessive whitespace
        content = re.sub(r'\s+', ' ', content)
        
        # Remove non-printable characters
        content = ''.join(char for char in content if char.isprintable() or char.isspace())
        
        # Normalize line breaks
        content = content.replace('\r\n', '\n').replace('\r', '\n')
        
        # Trim
        content = content.strip()
        
        return content
    
    def _create_chunks(self, content: str, file_type: DocumentType) -> List[DocumentChunk]:
        """Split content into chunks"""
        chunks = []
        
        # For structured documents, try to split by natural boundaries
        if file_type in [DocumentType.MARKDOWN, DocumentType.HTML]:
            # Split by headings
            sections = self._split_by_sections(content)
            for i, section in enumerate(sections):
                chunks.append(DocumentChunk(
                    text=section,
                    chunk_id=i,
                    metadata={'type': 'section'}
                ))
        else:
            # Split by approximate token count
            words = content.split()
            chunk_words = []
            chunk_id = 0
            
            for i, word in enumerate(words):
                chunk_words.append(word)
                
                # Check if we need to create a new chunk
                if len(chunk_words) >= self.chunk_size:
                    # Create chunk
                    chunks.append(DocumentChunk(
                        text=' '.join(chunk_words),
                        chunk_id=chunk_id,
                        metadata={'type': 'token_chunk'}
                    ))
                    chunk_id += 1
                    
                    # Keep overlap
                    chunk_words = chunk_words[-self.chunk_overlap:] if self.chunk_overlap > 0 else []
            
            # Add remaining words
            if chunk_words:
                chunks.append(DocumentChunk(
                    text=' '.join(chunk_words),
                    chunk_id=chunk_id,
                    metadata={'type': 'token_chunk'}
                ))
        
        return chunks
    
    def _split_by_sections(self, content: str) -> List[str]:
        """Split content by headings/sections"""
        sections = []
        current_section = ""
        
        lines = content.split('\n')
        for line in lines:
            line = line.strip()
            
            # Check for Markdown headings
            if line.startswith('#'):
                # Save current section
                if current_section:
                    sections.append(current_section)
                current_section = line + "\n"
            else:
                current_section += line + "\n" if line else ""
        
        # Add last section
        if current_section:
            sections.append(current_section)
        
        return sections
    
    def extract_sections(self, document: ProcessedDocument) -> List[Dict[str, Any]]:
        """Extract sections/headings from a document"""
        sections = []
        
        if document.metadata.file_type == DocumentType.MARKDOWN:
            # Parse Markdown headings
            lines = document.content.split('\n')
            current_level = 0
            current_section = None
            
            for line in lines:
                line = line.strip()
                
                # Check for heading
                if line.startswith('#'):
                    level = len(line.split()[0]) if line.split() else 0
                    title = line.lstrip('#').strip()
                    
                    # Save previous section
                    if current_section:
                        sections.append(current_section)
                    
                    current_section = {
                        'title': title,
                        'level': level,
                        'content': '',
                        'children': []
                    }
                    current_level = level
                elif line and current_section:
                    # Add to current section
                    current_section['content'] += line + "\n"
            
            # Add last section
            if current_section:
                sections.append(current_section)
        else:
            # For non-markdown, create sections based on chunks
            for i, chunk in enumerate(document.chunks):
                sections.append({
                    'title': f"Section {i+1}",
                    'level': 1,
                    'content': chunk.text,
                    'children': []
                })
        
        return sections
    
    def cleanup(self) -> None:
        """Clean up temporary files"""
        for temp_file in self._temp_files:
            try:
                os.remove(temp_file)
            except:
                pass
        self._temp_files.clear()
    
    def get_supported_extensions(self) -> List[str]:
        """Get list of supported file extensions"""
        return [
            '.pdf', '.docx', '.txt', '.md', '.markdown',
            '.html', '.htm', '.csv', '.xlsx', '.xls',
            '.pptx', '.ppt'
        ]
